from flask import Flask, request, send_file, jsonify
from pptx import Presentation
import copy
import io
import uuid
import requests as req_lib

try:
    from pypdf import PdfWriter, PdfReader
except ImportError:
    from PyPDF2 import PdfWriter, PdfReader

app = Flask(__name__)

# In-memory storage for merged files (demo purposes)
merged_files = {}

@app.route('/merge', methods=['POST'])
def merge():
    data = request.get_json()
    if not data or 'urls' not in data:
        return jsonify({'error': 'urls alanı eksik'}), 400

    urls = data['urls']
    if not urls:
        return jsonify({'error': 'Dosya URL listesi boş'}), 400

    pptx_buffers = []
    pdf_buffers = []

    for url in urls:
        try:
            response = req_lib.get(url, timeout=30)
            response.raise_for_status()
            buf = io.BytesIO(response.content)
            url_lower = url.lower().split('?')[0]
            if url_lower.endswith('.pptx'):
                pptx_buffers.append(buf)
            elif url_lower.endswith('.pdf'):
                pdf_buffers.append(buf)
        except Exception as e:
            return jsonify({'error': f'Dosya indirilemedi: {url} — {str(e)}'}), 400

    if not pptx_buffers and not pdf_buffers:
        return jsonify({'error': 'Desteklenmeyen dosya türü (sadece .pptx ve .pdf)'}), 400

    file_id = str(uuid.uuid4())

    if pptx_buffers:
        merged_pptx = Presentation()
        for buf in pptx_buffers:
            prs = Presentation(buf)
            for slide in prs.slides:
                layout = merged_pptx.slide_layouts[6]
                new_slide = merged_pptx.slides.add_slide(layout)
                for shape in slide.shapes:
                    el = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert(2, el)
        output = io.BytesIO()
        merged_pptx.save(output)
        output.seek(0)
        merged_files[file_id] = {'data': output, 'filename': 'merged.pptx', 'mimetype': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}

    elif pdf_buffers:
        pdf_writer = PdfWriter()
        for buf in pdf_buffers:
            pdf_reader = PdfReader(buf)
            for page in pdf_reader.pages:
                pdf_writer.add_page(page)
        output = io.BytesIO()
        pdf_writer.write(output)
        output.seek(0)
        merged_files[file_id] = {'data': output, 'filename': 'merged.pdf', 'mimetype': 'application/pdf'}

    download_url = f'https://sunum-birlestirici.onrender.com/download/{file_id}'
    return download_url, 200, {'Content-Type': 'text/plain'}


@app.route('/download/<file_id>', methods=['GET'])
def download(file_id):
    if file_id not in merged_files:
        return jsonify({'error': 'Dosya bulunamadı veya süresi doldu'}), 404
    f = merged_files[file_id]
    f['data'].seek(0)
    return send_file(f['data'], download_name=f['filename'], as_attachment=True, mimetype=f['mimetype'])


@app.route('/', methods=['GET'])
def index():
    return '''<!DOCTYPE html>
<html lang="tr">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Sunum Birleştirici</title>
  <style>
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #f5f5f5; display: flex; justify-content: center; align-items: center; min-height: 100vh; }
    .card { background: white; border-radius: 12px; padding: 48px; max-width: 480px; width: 100%; box-shadow: 0 2px 16px rgba(0,0,0,0.08); text-align: center; }
    h1 { font-size: 22px; font-weight: 600; color: #1a1a1a; margin-bottom: 8px; }
    p { color: #666; font-size: 14px; margin-bottom: 32px; line-height: 1.5; }
    .badge { display: inline-block; background: #e8f5e9; color: #2e7d32; font-size: 12px; font-weight: 500; padding: 4px 12px; border-radius: 20px; margin-bottom: 24px; }
    .formats { display: flex; justify-content: center; gap: 12px; }
    .fmt { background: #f0f4ff; color: #3b5bdb; border-radius: 8px; padding: 12px 24px; font-size: 13px; font-weight: 500; }
  </style>
</head>
<body>
  <div class="card">
    <span class="badge">&#x2713; Çevrimiçi</span>
    <h1>Sunum Birleştirici</h1>
    <p>Bu servis NocoBase ile entegre çalışır.<br>Dosyalarınızı NocoBase üzerinden yükleyip birleştirebilirsiniz.</p>
    <div class="formats">
      <div class="fmt">PPTX</div>
      <div class="fmt">PDF</div>
    </div>
  </div>
</body>
</html>''', 200, {'Content-Type': 'text/html'}


@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'})


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)
